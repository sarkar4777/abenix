"""PowerPoint/presentation analysis - extract slides, text, images, speaker notes, layout."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from engine.tools.base import BaseTool, ToolResult


class PresentationAnalyzerTool(BaseTool):
    name = "presentation_analyzer"
    description = (
        "Analyze PowerPoint presentations (.pptx): extract slide content, speaker notes, "
        "images, tables, charts, slide layouts, and master slides. Provides structured "
        "overview of presentation flow, content density per slide, and text extraction."
    )
    input_schema: dict[str, Any] = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to the PowerPoint file (.pptx)",
            },
            "operation": {
                "type": "string",
                "enum": ["overview", "slide", "all_text", "notes", "tables", "search"],
                "description": "Analysis operation",
                "default": "overview",
            },
            "slide_number": {
                "type": "integer",
                "description": "Specific slide number to analyze (1-indexed)",
            },
            "search_term": {
                "type": "string",
                "description": "Text to search across slides",
            },
        },
        "required": [],
    }

    async def execute(self, arguments: dict[str, Any]) -> ToolResult:
        from engine.tools._inline_file import materialise_path

        operation = arguments.get("operation", "overview")
        norm = dict(arguments)
        if arguments.get("file_path") and not norm.get("path"):
            norm["path"] = arguments["file_path"]
        resolved, err = materialise_path(norm, default_ext="pptx")
        if err:
            return ToolResult(content=err, is_error=True)
        file_path = resolved

        path = Path(file_path)
        if not path.exists():
            return ToolResult(content=f"File not found: {file_path}", is_error=True)

        if path.suffix.lower() not in (".pptx",):
            return ToolResult(
                content=f"Unsupported format: {path.suffix}. Only .pptx is supported.",
                is_error=True,
            )

        try:
            from pptx import Presentation

            prs = Presentation(str(path))

            if operation == "overview":
                result = self._overview(prs, path)
            elif operation == "slide":
                result = self._slide_detail(prs, arguments.get("slide_number", 1))
            elif operation == "all_text":
                result = self._all_text(prs)
            elif operation == "notes":
                result = self._all_notes(prs)
            elif operation == "tables":
                result = self._all_tables(prs)
            elif operation == "search":
                result = self._search(prs, arguments.get("search_term", ""))
            else:
                result = self._overview(prs, path)

            output = json.dumps(result, indent=2, default=str)
            return ToolResult(
                content=output, metadata={"file": str(path), "operation": operation}
            )
        except ImportError:
            return ToolResult(
                content="python-pptx library is required. Install with: pip install python-pptx",
                is_error=True,
            )
        except Exception as e:
            return ToolResult(
                content=f"Presentation analysis error: {e}", is_error=True
            )

    def _overview(self, prs: Any, path: Path) -> dict[str, Any]:
        slides_info = []
        total_text_chars = 0
        total_images = 0
        total_tables = 0

        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            image_count = 0
            table_count = 0
            has_chart = False

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
                if shape.shape_type == 13:
                    image_count += 1
                if shape.has_table:
                    table_count += 1
                if shape.has_chart:
                    has_chart = True

            text_content = " ".join(slide_text)
            total_text_chars += len(text_content)
            total_images += image_count
            total_tables += table_count

            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            slides_info.append(
                {
                    "slide": i,
                    "title": slide_text[0] if slide_text else "(no title)",
                    "text_preview": text_content[:200],
                    "word_count": len(text_content.split()),
                    "images": image_count,
                    "tables": table_count,
                    "has_chart": has_chart,
                    "has_notes": bool(notes),
                    "layout": slide.slide_layout.name if slide.slide_layout else "",
                }
            )

        return {
            "file": str(path),
            "slide_count": len(prs.slides),
            "total_characters": total_text_chars,
            "total_images": total_images,
            "total_tables": total_tables,
            "presentation_dimensions": {
                "width_inches": round(prs.slide_width.inches, 2),
                "height_inches": round(prs.slide_height.inches, 2),
            },
            "slides": slides_info,
        }

    def _slide_detail(self, prs: Any, slide_num: int) -> dict[str, Any]:
        if slide_num < 1 or slide_num > len(prs.slides):
            return {
                "error": f"Slide {slide_num} not found. Total slides: {len(prs.slides)}"
            }

        slide = prs.slides[slide_num - 1]
        shapes = []

        for shape in slide.shapes:
            shape_info: dict[str, Any] = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "position": {
                    "left": round(shape.left.inches, 2) if shape.left else 0,
                    "top": round(shape.top.inches, 2) if shape.top else 0,
                    "width": round(shape.width.inches, 2) if shape.width else 0,
                    "height": round(shape.height.inches, 2) if shape.height else 0,
                },
            }

            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        paragraphs.append(
                            {
                                "text": para.text,
                                "level": para.level,
                                "bold": any(
                                    run.font.bold
                                    for run in para.runs
                                    if run.font.bold is not None
                                ),
                            }
                        )
                shape_info["text"] = paragraphs

            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text for cell in row.cells])
                shape_info["table"] = table_data

            if shape.has_chart:
                shape_info["has_chart"] = True

            shapes.append(shape_info)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        return {
            "slide_number": slide_num,
            "layout": slide.slide_layout.name if slide.slide_layout else "",
            "shape_count": len(shapes),
            "shapes": shapes,
            "notes": notes,
        }

    def _all_text(self, prs: Any) -> dict[str, Any]:
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text.strip())
            slides_text.append(
                {
                    "slide": i,
                    "text": "\n".join(text_parts),
                }
            )

        full_text = "\n\n".join(
            f"--- Slide {s['slide']} ---\n{s['text']}" for s in slides_text if s["text"]
        )
        return {
            "total_slides": len(prs.slides),
            "slides_with_text": sum(1 for s in slides_text if s["text"]),
            "full_text": full_text,
            "total_words": len(full_text.split()),
        }

    def _all_notes(self, prs: Any) -> dict[str, Any]:
        notes = []
        for i, slide in enumerate(prs.slides, 1):
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note_text = slide.notes_slide.notes_text_frame.text.strip()
                if note_text:
                    notes.append({"slide": i, "notes": note_text})

        return {"slides_with_notes": len(notes), "notes": notes}

    def _all_tables(self, prs: Any) -> dict[str, Any]:
        tables = []
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text for cell in row.cells])
                    tables.append(
                        {
                            "slide": i,
                            "rows": len(rows),
                            "columns": len(rows[0]) if rows else 0,
                            "data": rows,
                        }
                    )

        return {"table_count": len(tables), "tables": tables}

    def _search(self, prs: Any, term: str) -> dict[str, Any]:
        if not term:
            return {"error": "search_term is required"}

        matches = []
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if term.lower() in para.text.lower():
                            matches.append(
                                {
                                    "slide": i,
                                    "shape": shape.name,
                                    "text": para.text,
                                }
                            )

        return {
            "search_term": term,
            "match_count": len(matches),
            "matches": matches[:100],
        }
